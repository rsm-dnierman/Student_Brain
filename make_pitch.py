#!/usr/bin/env python3
"""
MSBA Brain — Demo Pitch Deck (python-pptx)
Run from project root:  python make_pitch.py
Output:                 MSBA_Brain_Pitch.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
import os

# ── Palette ───────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0D, 0x11, 0x17)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT    = RGBColor(0xE6, 0xED, 0xF3)
MUTED    = RGBColor(0x7D, 0x8B, 0x99)
ACCENT   = RGBColor(0x1F, 0x6F, 0xEB)   # primary blue #1f6feb
ACCENT_L = RGBColor(0x58, 0xA6, 0xFF)   # lighter blue
CARD_BG  = RGBColor(0x16, 0x1B, 0x22)
GREEN    = RGBColor(0x2E, 0xA0, 0x43)
YELLOW   = RGBColor(0xE3, 0xB3, 0x41)
BORDER   = RGBColor(0x30, 0x3C, 0x50)

W, H = 13.333, 7.5   # 16:9 widescreen

# ── Paths ─────────────────────────────────────────────────────────────────────
BASE   = os.path.dirname(os.path.abspath(__file__))
SS_DIR = os.path.join(BASE, "screenshots")

IMG = {
    "sources"   : os.path.join(SS_DIR, "Screenshot 2026-06-02 at 4.19.43 PM.png"),
    "preview"   : os.path.join(SS_DIR, "Screenshot 2026-06-02 at 4.19.59 PM.png"),
    "setup"     : os.path.join(SS_DIR, "Screenshot 2026-06-02 at 4.20.19 PM.png"),
    "chat"      : os.path.join(SS_DIR, "Screenshot 2026-06-02 at 4.21.13 PM.png"),
    "citations" : os.path.join(SS_DIR, "Screenshot 2026-06-02 at 4.21.42 PM.png"),
    "flashcards": os.path.join(SS_DIR, "Screenshot 2026-06-02 at 4.23.13 PM.png"),
}

# ── Core helpers ──────────────────────────────────────────────────────────────
def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    f = s.background.fill
    f.solid()
    f.fore_color.rgb = BG
    return s

def make_tf(slide, x, y, w, h):
    sh = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    sh.text_frame.word_wrap = True
    return sh.text_frame

def ap(tf_obj, text, size=16, color=LIGHT, bold=False, italic=False,
       align=PP_ALIGN.LEFT, sb=0, sa=0):
    """Add a paragraph to a text frame; reuses the first empty paragraph."""
    pars = tf_obj.paragraphs
    if len(pars) == 1 and not pars[0].runs:
        p = pars[0]
    else:
        p = tf_obj.add_paragraph()
    p.alignment = align
    if sb: p.space_before = Pt(sb)
    if sa: p.space_after = Pt(sa)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"

def hbar(slide, x, y, w, color=ACCENT, h=0.04):
    s = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()

def add_rect(slide, x, y, w, h, fill=CARD_BG, border=None, rounded=True):
    t = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    s = slide.shapes.add_shape(t, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s

def add_pic(slide, key, x, y, w=None, h=None):
    kw = {}
    if w: kw["width"]  = Inches(w)
    if h: kw["height"] = Inches(h)
    return slide.shapes.add_picture(IMG[key], Inches(x), Inches(y), **kw)

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def stripes(slide):
    hbar(slide, 0, 0,        W, color=ACCENT, h=0.055)
    hbar(slide, 0, H - 0.055, W, color=ACCENT, h=0.055)


# ── Slide 1: Title ────────────────────────────────────────────────────────────
def slide_title(prs):
    s = new_slide(prs)
    stripes(s)

    t = make_tf(s, 0.7, 0.85, 7.2, 0.4)
    ap(t, "RADY SCHOOL OF MANAGEMENT  ·  UCSD  ·  2026", 10, MUTED)

    t = make_tf(s, 0.7, 1.35, 7.2, 2.6)
    ap(t, "Your Semester of Notes,", 44, WHITE, bold=True)
    ap(t, "Now Searchable.", 44, ACCENT_L, bold=True)

    hbar(s, 0.7, 3.85, 5.0)

    t = make_tf(s, 0.7, 4.05, 7.0, 0.65)
    ap(t, "A personal knowledge base for MSBA course material", 17, MUTED, italic=True)

    t = make_tf(s, 0.7, 5.0, 5.0, 0.8)
    ap(t, "Daniel Nierman", 15, LIGHT, bold=True)
    ap(t, "MS Business Analytics  ·  Rady School of Management", 12, MUTED)

    # Right: chat screenshot — the most compelling view of the finished product
    add_pic(s, "chat", 8.0, 0.65, w=5.1)

    add_notes(s, "Hi everyone. I'm going to show you something I built during the program to solve "
                  "a problem I know every one of us has had at some point.")
    return s


# ── Slide 2: The Problem ──────────────────────────────────────────────────────
def slide_problem(prs):
    s = new_slide(prs)
    stripes(s)

    t = make_tf(s, 0.6, 0.3, 5.0, 0.35)
    ap(t, "THE PROBLEM", 10, ACCENT, bold=True)

    t = make_tf(s, 0.6, 0.65, 9.5, 0.8)
    ap(t, "Sound familiar?", 34, WHITE, bold=True)

    # Quote box
    add_rect(s, 0.6, 1.55, 12.1, 1.55, fill=RGBColor(0x0F, 0x18, 0x28), border=ACCENT, rounded=False)
    t = make_tf(s, 0.95, 1.7, 11.5, 1.25)
    ap(t, "\"It's midnight before finals. I know there was a slide about regularization "
          "somewhere — was it in the ML course or the analytics one? Let me just check Canvas real quick.\"",
       17, LIGHT, italic=True)

    # Three pain-point cards
    pain = [
        ("Canvas is a maze",
         "Modules buried under weeks of content. Click, scroll, download, repeat. "
         "Every file requires its own hunt."),
        ("100+ files, zero unified search",
         "PDFs, notebooks, slides scattered across 8 courses. "
         "No way to search across all of them at once."),
        ("By finals week: total chaos",
         "You're not sure what you have, what you've read, "
         "or where anything lives anymore."),
    ]

    for i, (title, desc) in enumerate(pain):
        cx = 0.6 + i * 4.15
        add_rect(s, cx, 3.3, 3.9, 3.0, fill=CARD_BG, border=BORDER)
        t2 = make_tf(s, cx + 0.22, 3.5, 3.5, 0.5)
        ap(t2, title, 15, ACCENT_L, bold=True)
        hbar(s, cx + 0.22, 3.98, 3.3, color=ACCENT, h=0.025)
        t3 = make_tf(s, cx + 0.22, 4.1, 3.5, 2.0)
        ap(t3, desc, 13, LIGHT)

    add_notes(s, "Canvas is great for submitting homework. Not so great for studying from your material "
                  "six months later. By the end of the program you have hundreds of files and no way to find anything quickly.")
    return s


# ── Slide 3: Solution Overview ────────────────────────────────────────────────
def slide_solution(prs):
    s = new_slide(prs)
    stripes(s)

    t = make_tf(s, 0.6, 0.3, 8.0, 0.35)
    ap(t, "THE SOLUTION", 10, ACCENT, bold=True)

    t = make_tf(s, 0.6, 0.65, 10.5, 0.8)
    ap(t, "Four steps from scattered files to searchable brain", 32, WHITE, bold=True)

    hbar(s, 0.6, 1.45, 12.1)

    steps = [
        ("01", "SCRAPE",   ACCENT,   "Connect Canvas, RSM Django sites, and any URL. One click downloads everything to your machine."),
        ("02", "ORGANIZE", GREEN,    "Files sorted by course. PDFs, notebooks, and slides parsed into clean, structured text."),
        ("03", "INDEX",    YELLOW,   "Content chunked, embedded locally with sentence-transformers, stored in ChromaDB. Fully offline."),
        ("04", "QUERY",    ACCENT_L, "Ask in plain English. Get cited answers with a direct link to open the original source file."),
    ]

    for i, (num, title, color, desc) in enumerate(steps):
        cx = 0.5 + i * 3.2

        add_rect(s, cx, 1.7, 2.95, 4.65, fill=CARD_BG, border=color)

        t2 = make_tf(s, cx + 0.22, 1.92, 2.55, 0.65)
        ap(t2, num, 30, color, bold=True)

        t3 = make_tf(s, cx + 0.22, 2.55, 2.55, 0.5)
        ap(t3, title, 14, WHITE, bold=True)

        hbar(s, cx + 0.22, 3.03, 2.35, color=color, h=0.025)

        t4 = make_tf(s, cx + 0.22, 3.15, 2.52, 2.9)
        ap(t4, desc, 13, LIGHT)

        if i < 3:
            t5 = make_tf(s, cx + 3.02, 3.8, 0.3, 0.4)
            ap(t5, "→", 20, ACCENT_L, bold=True)

    add_notes(s, "Four tabs in the app, four steps in the pipeline. "
                  "Everything is local — no data leaves your machine except the final query to Claude.")
    return s


# ── Slide 4: Collect (Steps 1–2) ─────────────────────────────────────────────
def slide_collect(prs):
    s = new_slide(prs)
    stripes(s)

    t = make_tf(s, 0.6, 0.3, 8.0, 0.35)
    ap(t, "STEPS 1–2: CONNECT & COLLECT", 10, ACCENT, bold=True)

    t = make_tf(s, 0.6, 0.65, 7.5, 0.8)
    ap(t, "Point it at your courses. Let it do the work.", 28, WHITE, bold=True)

    hbar(s, 0.6, 1.42, 12.1)

    # Right: preview screenshot (charts make it the most compelling visual)
    add_pic(s, "preview", 7.0, 1.55, w=6.1)

    # Left: what it connects to
    t = make_tf(s, 0.6, 1.65, 6.1, 0.4)
    ap(t, "Three source types:", 13, ACCENT_L, bold=True)

    connects = [
        ("Canvas (pre-loaded)",  "Auto-connects via your personal API token. All course files."),
        ("Rady Django sites",    "RSM course sites at rsm-django-02.ucsd.edu. Login once, scrapes everything."),
        ("Any URL",              "Drop in a lecture page, external resource, or course website."),
    ]

    cy = 2.15
    for src, desc in connects:
        add_rect(s, 0.6, cy, 6.1, 0.92, fill=CARD_BG, border=BORDER, rounded=False)
        t2 = make_tf(s, 0.82, cy + 0.1, 5.7, 0.32)
        ap(t2, src, 13, ACCENT_L, bold=True)
        t3 = make_tf(s, 0.82, cy + 0.42, 5.7, 0.42)
        ap(t3, desc, 12, LIGHT)
        cy += 1.05

    # Stat row
    stats = [("110", "files"), ("8", "courses"), ("541 MB", "downloaded")]
    for i, (num, label) in enumerate(stats):
        cx = 0.6 + i * 2.05
        add_rect(s, cx, 5.45, 1.9, 1.1, fill=CARD_BG, border=ACCENT)
        t4 = make_tf(s, cx + 0.1, 5.56, 1.7, 0.48)
        ap(t4, num, 22, ACCENT_L, bold=True, align=PP_ALIGN.CENTER)
        t5 = make_tf(s, cx + 0.1, 6.0, 1.7, 0.3)
        ap(t5, label, 11, MUTED, align=PP_ALIGN.CENTER)

    add_notes(s, "Step 1 is data sources — Canvas is pre-configured once you paste your API token. "
                  "RSM Django sites take your UCSD login. Step 2 shows you a live preview before indexing: "
                  "110 files across 8 courses, 541 MB total, broken down by course and file type.")
    return s


# ── Slide 5: Index (Step 3) ───────────────────────────────────────────────────
def slide_index(prs):
    s = new_slide(prs)
    stripes(s)

    t = make_tf(s, 0.6, 0.3, 8.0, 0.35)
    ap(t, "STEP 3: BUILD THE BRAIN", 10, ACCENT, bold=True)

    t = make_tf(s, 0.6, 0.65, 8.5, 0.8)
    ap(t, "From files to a searchable knowledge base", 28, WHITE, bold=True)

    hbar(s, 0.6, 1.42, 12.1)

    # Left: AI setup screenshot
    add_pic(s, "setup", 0.5, 1.6, w=6.3)

    # Right: four process steps
    t2 = make_tf(s, 7.2, 1.65, 5.8, 0.4)
    ap(t2, "What happens when you hit Index:", 13, ACCENT_L, bold=True)

    process = [
        ("PARSE",  "PDFs, notebooks, slides → clean text. PyMuPDF handles the heavy lifting; Claude Vision as fallback for image-heavy PDFs."),
        ("CHUNK",  "Text split into ~500-token overlapping chunks so nothing gets cut off mid-concept."),
        ("EMBED",  "Each chunk encoded with sentence-transformers/all-MiniLM-L6-v2. Local model — no data leaves your machine."),
        ("STORE",  "Chunks + vectors saved to ChromaDB. 1,015 chunks across 8 courses, ready for hybrid search."),
    ]

    cy = 2.15
    for step_label, step_desc in process:
        add_rect(s, 7.2, cy, 5.9, 1.05, fill=CARD_BG, border=BORDER, rounded=False)
        t3 = make_tf(s, 7.42, cy + 0.1, 5.5, 0.3)
        ap(t3, step_label, 10, ACCENT_L, bold=True)
        t4 = make_tf(s, 7.42, cy + 0.4, 5.5, 0.55)
        ap(t4, step_desc, 12, LIGHT)
        cy += 1.15

    add_notes(s, "The indexing step is fully automated. The only thing you supply is your Anthropic API key "
                  "and a model choice. Everything else — parsing, chunking, embedding, storing — happens locally. "
                  "The sentence-transformer model is downloaded once and runs on your CPU.")
    return s


# ── Slide 6: Query (Step 4) ───────────────────────────────────────────────────
def slide_query(prs):
    s = new_slide(prs)
    stripes(s)

    t = make_tf(s, 0.6, 0.3, 9.0, 0.35)
    ap(t, "STEP 4: ASK ANYTHING", 10, ACCENT, bold=True)

    t = make_tf(s, 0.6, 0.65, 9.0, 0.8)
    ap(t, "Natural language in. Cited answers out.", 28, WHITE, bold=True)

    hbar(s, 0.6, 1.42, 12.1)

    # Side-by-side screenshots
    add_pic(s, "chat",      0.4,  1.6, w=6.3)
    add_pic(s, "citations", 6.85, 1.6, w=6.3)

    # Captions
    t2 = make_tf(s, 0.4, 6.12, 6.3, 0.45)
    ap(t2, "Ask scoped to one course or across all 8 — your choice", 11, MUTED, italic=True, align=PP_ALIGN.CENTER)

    t3 = make_tf(s, 6.85, 6.12, 6.3, 0.45)
    ap(t3, "Each source shows relevance score + Open File button → downloads the original PDF", 11, MUTED, italic=True, align=PP_ALIGN.CENTER)

    add_notes(s, "The retrieval system runs three passes: BM25 keyword search, vector similarity, "
                  "and a cross-encoder reranker to pick the best chunks. Claude gets those chunks as context "
                  "and generates a cited answer. Every claim is traceable back to a specific file and page.")
    return s


# ── Slide 7: Live Demo ────────────────────────────────────────────────────────
def slide_demo(prs):
    s = new_slide(prs)
    stripes(s)

    t = make_tf(s, 0.6, 0.3, 6.0, 0.35)
    ap(t, "LIVE DEMO", 10, GREEN, bold=True)

    t = make_tf(s, 0.6, 0.65, 6.5, 0.8)
    ap(t, "Let's see it.", 32, WHITE, bold=True)

    hbar(s, 0.6, 1.42, 5.8)

    demo_steps = [
        "Add Canvas as a data source (pre-loaded, just show the token config)",
        "Scrape & Preview — 110 files, 8 courses, 541 MB",
        "AI Setup — 1,015 chunks indexed and ready",
        "Ask: \"Explain the evaluation criteria for clustering solutions\"",
        "Walk through the cited answer — 4 sources, relevance scores",
        "Open source PDF directly from the app",
        "Bonus: Study Tools → generate flashcards on A/B testing",
    ]

    t2 = make_tf(s, 0.6, 1.6, 6.3, 5.5)
    for i, step in enumerate(demo_steps):
        ap(t2, f"  {i+1}.  {step}", 14, LIGHT, sb=5)

    # Right: flashcards screenshot
    add_pic(s, "flashcards", 7.1, 1.55, w=6.0)

    add_notes(s, "Walk through each step live. The Study Tools tab is a bonus — it's not just a search engine, "
                  "it actively helps you study.")
    return s


# ── Slide 8: Value Add ────────────────────────────────────────────────────────
def slide_value(prs):
    s = new_slide(prs)
    stripes(s)

    t = make_tf(s, 0.6, 0.3, 9.0, 0.35)
    ap(t, "VALUE ADD", 10, ACCENT, bold=True)

    t = make_tf(s, 0.6, 0.65, 9.0, 0.8)
    ap(t, "So what does this actually give you?", 30, WHITE, bold=True)

    hbar(s, 0.6, 1.42, 12.1)

    cards = [
        (ACCENT,   "Time Saved",
         "Find any concept in seconds. No more hunting through Canvas modules "
         "or Ctrl+F-ing a 200-slide deck at midnight before finals."),
        (GREEN,    "Better Retention",
         "Study Tools generates flashcards and quizzes on demand from your actual course material. "
         "Active recall beats passive rereading every time."),
        (YELLOW,   "Yours to Keep",
         "The database is local. When Canvas disappears after graduation, "
         "your notes — and your brain — don't go with it."),
        (ACCENT_L, "Transferable Skill",
         "Same architecture works on any document collection: new job onboarding docs, "
         "client materials, research papers. Build a brain for anything."),
    ]

    for i, (color, title, desc) in enumerate(cards):
        cx = 0.5 + i * 3.2
        add_rect(s, cx, 1.7, 2.95, 4.7, fill=CARD_BG, border=color)
        t2 = make_tf(s, cx + 0.22, 1.92, 2.55, 0.5)
        ap(t2, title, 15, WHITE, bold=True)
        hbar(s, cx + 0.22, 2.42, 2.35, color=color, h=0.025)
        t3 = make_tf(s, cx + 0.22, 2.55, 2.52, 3.6)
        ap(t3, desc, 13, LIGHT)

    add_notes(s, "Three takeaways: speed, retention, and portability. "
                  "The portability piece is the one people don't think about until Canvas access expires.")
    return s


# ── Slide 9: Close ────────────────────────────────────────────────────────────
def slide_close(prs):
    s = new_slide(prs)
    stripes(s)

    t = make_tf(s, 1.2, 1.1, 10.9, 2.6)
    ap(t, "Built it because I needed it.", 42, WHITE, bold=True, align=PP_ALIGN.CENTER)
    ap(t, "Turns out, so did everyone else.", 42, ACCENT_L, bold=True, align=PP_ALIGN.CENTER)

    hbar(s, 3.2, 3.6, 6.9, color=MUTED, h=0.02)

    t2 = make_tf(s, 2.0, 3.8, 9.3, 0.65)
    ap(t2, "All 8 MSBA courses. 110 files. 1,015 chunks. One search bar.",
       17, MUTED, italic=True, align=PP_ALIGN.CENTER)

    # Stats row
    stats = [("110", "Files Indexed"), ("1,015", "Searchable Chunks"), ("8", "MSBA Courses")]
    for i, (num, label) in enumerate(stats):
        cx = 2.3 + i * 3.15
        add_rect(s, cx, 4.6, 2.85, 1.55, fill=CARD_BG, border=ACCENT)
        t3 = make_tf(s, cx + 0.1, 4.72, 2.65, 0.7)
        ap(t3, num, 34, ACCENT_L, bold=True, align=PP_ALIGN.CENTER)
        t4 = make_tf(s, cx + 0.1, 5.38, 2.65, 0.45)
        ap(t4, label, 12, MUTED, align=PP_ALIGN.CENTER)

    t5 = make_tf(s, 3.5, 6.45, 6.3, 0.5)
    ap(t5, "Questions?", 22, WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_notes(s, "Thanks. Happy to demo any specific part in more depth, or talk through the architecture. "
                  "Code is on GitHub if anyone wants to run it on their own courses.")
    return s


# ── Main ──────────────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)

    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_collect(prs)
    slide_index(prs)
    slide_query(prs)
    slide_demo(prs)
    slide_value(prs)
    slide_close(prs)

    out = os.path.join(BASE, "MSBA_Brain_Pitch.pptx")
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
